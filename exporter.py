import os
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches, Pt
import markdown2

class ReportExporter:
    @staticmethod
    def to_markdown(content: str, filename: str = "report.md"):
        with open(filename, "w", encoding="utf-8") as f:
            f.write(content)
        return filename

    @staticmethod
    def to_pdf(content_md: str, filename: str = "report.pdf"):
        # Convert MD to HTML for better parsing if needed, but SimpleDocTemplate is simpler
        styles = getSampleStyleSheet()
        doc = SimpleDocTemplate(filename, pagesize=letter)
        story = []
        
        # Split content by lines and attempt basic markdown-like parsing
        for line in content_md.split('\n'):
            if line.startswith('# '):
                story.append(Paragraph(line[2:], styles['Title']))
            elif line.startswith('## '):
                story.append(Paragraph(line[3:], styles['Heading2']))
            elif line.strip() == "":
                story.append(Spacer(1, 12))
            else:
                story.append(Paragraph(line, styles['Normal']))
        
        doc.build(story)
        return filename

    @staticmethod
    def to_slides(topic: str, content: str, filename: str = "presentation.pptx"):
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = topic
        subtitle.text = "Autonomous Research Analyst Report"
        
        # Content Slides (Splitting content by double newline or headers)
        sections = content.split('## ')
        for section in sections[1:]:
            lines = section.split('\n')
            header = lines[0]
            body = "\n".join(lines[1:])
            
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = header
            tf = body_shape.text_frame
            tf.text = body[:1000] # Limit content per slide
            
        prs.save(filename)
        return filename
